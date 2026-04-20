"""Build the STAT 4830 presentation deck (python-pptx).

13 slides, 10–15 minutes, smart grad audience.
Three tracks: cascade GPU (1.4), Lasserre SDP (1.3 certified), MATLAB audit.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
FIG = HERE / "figures"
OUT = HERE / "STAT 4830 Presentation.pptx"

# --- palette ---
NAVY = RGBColor(0x0B, 0x3D, 0x91)
TEAL = RGBColor(0x13, 0x80, 0x86)
GOLD = RGBColor(0xD4, 0xA0, 0x17)
CORAL = RGBColor(0xC4, 0x45, 0x36)
GREY = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF5, 0xF6, 0xF8)
INK = RGBColor(0x11, 0x18, 0x27)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_prs() -> Presentation:
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, italic=False, color=INK,
             align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        lines = [text]
    else:
        lines = list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_rich(slide, left, top, width, height, runs, *,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs = list of paragraphs; each paragraph = list of (text, dict)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for text, style in para:
            r = p.add_run()
            r.text = text
            r.font.name = style.get("font", "Calibri")
            r.font.size = Pt(style.get("size", 18))
            r.font.bold = style.get("bold", False)
            r.font.italic = style.get("italic", False)
            r.font.color.rgb = style.get("color", INK)
    return tb


def add_rule(slide, left, top, width, color=NAVY, thickness=2.5):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(thickness)
    return line


def add_box(slide, left, top, width, height, fill=LIGHT, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def slide_header(slide, title, subtitle=None, color=NAVY):
    add_text(slide, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.7),
             title, size=30, bold=True, color=color, font="Calibri")
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.95), Inches(12.1), Inches(0.5),
                 subtitle, size=16, color=GREY, font="Calibri")
    add_rule(slide, Inches(0.6), Inches(1.45), Inches(12.1), color=color)


def add_footer(slide, idx, total):
    add_text(slide, Inches(0.6), Inches(7.1), Inches(8),
             Inches(0.3),
             "Piterbarg · Bajaj · Vincent   |   STAT 4830",
             size=10, color=GREY)
    add_text(slide, Inches(11.8), Inches(7.1), Inches(1.0),
             Inches(0.3), f"{idx} / {total}",
             size=10, color=GREY, align=PP_ALIGN.RIGHT)


# ============================================================
# SLIDES
# ============================================================

def slide_01_title(prs):
    s = blank_slide(prs)

    add_box(s, Inches(0), Inches(2.0), SLIDE_W, Inches(0.04), fill=NAVY)

    add_text(s, Inches(0.8), Inches(2.3), Inches(11.8), Inches(1.0),
             "Lower Bounds for the",
             size=40, bold=True, color=INK)
    add_text(s, Inches(0.8), Inches(2.95), Inches(11.8), Inches(1.0),
             "Autoconvolution Constant",
             size=40, bold=True, color=INK)
    add_text(s, Inches(0.8), Inches(3.85), Inches(11.8), Inches(0.7),
             "Three parallel attacks on  C₁ₐ",
             size=22, color=TEAL)

    add_rule(s, Inches(0.8), Inches(4.8), Inches(4.0), color=NAVY, thickness=1.5)

    add_text(s, Inches(0.8), Inches(4.95), Inches(11.8), Inches(0.6),
             "Andrei Piterbarg · Jai Bajaj · Derrick Vincent",
             size=20, bold=True, color=INK)
    add_text(s, Inches(0.8), Inches(5.55), Inches(11.8), Inches(0.5),
             "STAT 4830  —  Spring 2026",
             size=16, color=GREY)


def slide_02_problem(prs, idx, total):
    s = blank_slide(prs)
    slide_header(s, "The problem",
                 "How small can the peak of a normalized self-convolution be?")

    add_rich(s, Inches(0.6), Inches(1.8), Inches(6.3), Inches(3.2), [
        [("Given  ", {"size": 19}),
         ("f : ℝ → ℝ₊",  {"size": 19, "italic": True, "color": NAVY}),
         ("  with  supp f ⊆ [−¼, ¼]", {"size": 19})],
        [("and  ∫ f = 1,  define", {"size": 19})],
        [("", {"size": 10})],
        [("C₁ₐ  =  inf_f  ‖ f ∗ f ‖∞",
          {"size": 32, "bold": True, "color": NAVY})],
        [("", {"size": 10})],
        [("• connected to Sidon / generalized Sidon set densities",
          {"size": 15, "color": GREY})],
        [("• best-known bounds have not moved since 2017",
          {"size": 15, "color": GREY})],
        [("• every improvement is a real contribution",
          {"size": 15, "color": GREY})],
    ])

    s.shapes.add_picture(str(FIG / "fig_autoconv_intuition.png"),
                         Inches(7.0), Inches(2.0), width=Inches(6.0))

    add_box(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.2),
            fill=LIGHT, line=NAVY)
    add_rich(s, Inches(0.8), Inches(5.75), Inches(11.7), Inches(1.0), [
        [("Best published bounds", {"size": 14, "color": GREY})],
        [("1.2802", {"size": 24, "bold": True, "color": NAVY}),
         ("  ≤  ", {"size": 24}),
         ("C₁ₐ", {"size": 24, "bold": True}),
         ("  ≤  ", {"size": 24}),
         ("1.5029", {"size": 24, "bold": True, "color": CORAL}),
         ("       — unchanged since 2017",
          {"size": 16, "color": GREY, "italic": True})],
    ])

    add_footer(s, idx, total)


def slide_03_prior_work(prs, idx, total):
    s = blank_slide(prs)
    slide_header(s, "Prior work: Cloninger–Steinerberger 2017",
                 "Branch-and-prune over a fine height grid  Bₙ,ₘ")

    add_rich(s, Inches(0.6), Inches(1.7), Inches(7.3), Inches(4.8), [
        [("Step 1 · Discretize", {"size": 17, "bold": True, "color": NAVY})],
        [("partition  [−¼, ¼]  into  d = 2n  bins;  approximate  ", {"size": 15}),
         ("f", {"size": 15, "italic": True}),
         ("  by a step function  ", {"size": 15}),
         ("g", {"size": 15, "italic": True}),
         ("  with heights in  (1/m) ℕ", {"size": 15})],
        [("", {"size": 6})],
        [("Step 2 · Windowed test values",
          {"size": 17, "bold": True, "color": NAVY})],
        [("for each window  W,  lower-bound  ‖f ∗ f‖∞  by a quadratic form in the bin heights",
          {"size": 15})],
        [("", {"size": 6})],
        [("Step 3 · Prune", {"size": 17, "bold": True, "color": NAVY})],
        [("rule out any height vector whose test value exceeds the target threshold",
          {"size": 15})],
        [("", {"size": 6})],
        [("Step 4 · Exhaust", {"size": 17, "bold": True, "color": NAVY})],
        [("if all grid points are pruned, the bound is proved",
          {"size": 15})],
    ])

    # right panel: result card
    add_box(s, Inches(8.3), Inches(1.9), Inches(4.4), Inches(4.4),
            fill=LIGHT, line=NAVY)
    add_text(s, Inches(8.5), Inches(2.1), Inches(4.0), Inches(0.5),
             "Published result", size=14, color=GREY)
    add_text(s, Inches(8.5), Inches(2.45), Inches(4.0), Inches(1.2),
             "C₁ₐ  ≥  1.2802", size=34, bold=True, color=NAVY)
    add_text(s, Inches(8.5), Inches(3.6), Inches(4.0), Inches(0.5),
             "at  n = 24  (d = 48),  m = 50",
             size=15, color=INK)
    add_text(s, Inches(8.5), Inches(4.05), Inches(4.0), Inches(0.5),
             "≈ 20,000 CPU-hours",
             size=15, color=INK)
    add_rule(s, Inches(8.5), Inches(4.85), Inches(4.0), color=GREY, thickness=1.0)
    add_text(s, Inches(8.5), Inches(4.95), Inches(4.0), Inches(1.3),
             "Correction term  2/m + 1/m²  controls the step-function approximation error.",
             size=13, color=GREY)

    add_footer(s, idx, total)


def slide_04_three_tracks(prs, idx, total):
    s = blank_slide(prs)
    slide_header(s, "Our three tracks",
                 "One problem; three independent attacks; one MATLAB audit along the way")

    s.shapes.add_picture(str(FIG / "fig_three_tracks.png"),
                         Inches(0.6), Inches(1.9), width=Inches(12.1))

    add_text(s, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.3),
             "We independently implemented a branch-and-prune cascade on GPU  (→ 1.4), "
             "built a Lasserre SDP certificate  (→ 1.3, certified),  and, while "
             "reproducing the C&S baseline, audited the MATLAB code the authors shared. "
             "All three are presented here.",
             size=15, color=INK)

    add_footer(s, idx, total)


def slide_05_cascade_method(prs, idx, total):
    s = blank_slide(prs)
    slide_header(s, "Track 1 · The cascade",
                 "Recursive refinement with novel pruning bounds, fused GPU kernels",
                 color=NAVY)

    add_rich(s, Inches(0.6), Inches(1.7), Inches(6.7), Inches(5.0), [
        [("Novel pruning bounds", {"size": 18, "bold": True, "color": NAVY})],
        [("• contiguous-block sums catch interior concentrations", {"size": 15})],
        [("• two-max enhanced bound: cross-term dominance", {"size": 15})],
        [("• edge-pair cross-terms  c₀ · c_{d−1}", {"size": 15})],
        [("each bound proved, then compiled into the pruner", {"size": 13, "italic": True, "color": GREY})],
        [("", {"size": 8})],
        [("Fused CUDA kernel", {"size": 18, "bold": True, "color": NAVY})],
        [("• everything on-GPU; no global-memory round trips", {"size": 15})],
        [("• early exit on partial convolution sums  (most compositions die cheap)",
          {"size": 15})],
        [("• two-stage freeze kernel at deep levels: perturbation bound clears 90%+ of frozen cases,",
          {"size": 15})],
        [("  avoiding warp divergence → 10×–100× speedup", {"size": 15})],
    ])

    s.shapes.add_picture(str(FIG / "fig_cascade_schematic.png"),
                         Inches(7.5), Inches(1.85), width=Inches(5.5))

    add_footer(s, idx, total)


def slide_06_cascade_result(prs, idx, total):
    s = blank_slide(prs)
    slide_header(s, "Track 1 · Result",
                 "Cascade pushes the lower bound to 1.4",
                 color=NAVY)

    add_box(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(4.5),
            fill=LIGHT, line=NAVY)
    add_text(s, Inches(0.8), Inches(1.95), Inches(5.6), Inches(0.5),
             "Outcome", size=14, color=GREY)
    add_text(s, Inches(0.8), Inches(2.35), Inches(5.6), Inches(1.4),
             "C₁ₐ  ≥  1.4", size=52, bold=True, color=NAVY)
    add_text(s, Inches(0.8), Inches(3.75), Inches(5.6), Inches(0.5),
             "at  (d, m) = (128, 20)",
             size=16, color=INK)
    add_text(s, Inches(0.8), Inches(4.2), Inches(5.6), Inches(0.5),
             "≈ 70 hours on a 128-core machine",
             size=16, color=INK)
    add_rule(s, Inches(0.8), Inches(5.05), Inches(5.6), color=GREY, thickness=1.0)
    add_rich(s, Inches(0.8), Inches(5.15), Inches(5.6), Inches(1.2), [
        [("First significant lower-bound improvement since 2017.",
          {"size": 14, "color": INK})],
        [("Gap to upper bound narrowed by  ~50%  (0.22 → 0.10).",
          {"size": 14, "color": INK})],
    ])

    add_rich(s, Inches(7.0), Inches(1.9), Inches(5.8), Inches(5.0), [
        [("What the cascade buys you", {"size": 18, "bold": True, "color": NAVY})],
        [("", {"size": 6})],
        [("• sharp in principle — exhaustive enumeration rules out every surviving point",
          {"size": 15})],
        [("• scales to  d = 128  via symmetry reduction + GPU pruning",
          {"size": 15})],
        [("", {"size": 6})],
        [("What it doesn't", {"size": 18, "bold": True, "color": NAVY})],
        [("", {"size": 6})],
        [("• exponential in  d  in the worst case",
          {"size": 15})],
        [("• correctness rides on many small inequalities — each needs a proof",
          {"size": 15})],
        [("• motivated us to build a second, independent certifier …",
          {"size": 15, "italic": True, "color": GREY})],
    ])

    add_footer(s, idx, total)


def slide_07_sdp_method(prs, idx, total):
    s = blank_slide(prs)
    slide_header(s, "Track 2 · Lasserre SDP",
                 "Continuous → polynomial → semidefinite  (no enumeration)",
                 color=TEAL)

    add_rich(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.6), [
        [("Step 1.  Reduce", {"size": 18, "bold": True, "color": TEAL}),
         ("   the infinite-dimensional problem on  ", {"size": 18}),
         ("𝓕", {"size": 18, "italic": True}),
         ("  to a polynomial minimization on the simplex:",
          {"size": 18})],
    ])
    add_text(s, Inches(0.6), Inches(2.35), Inches(12.1), Inches(0.6),
             "val(d)  =  minᵤ ∈ Δ_d   max_W   μᵀ M_W μ   ≤  C₁ₐ",
             size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_rich(s, Inches(0.6), Inches(3.15), Inches(12.1), Inches(0.6), [
        [("Step 2.  Relax", {"size": 18, "bold": True, "color": TEAL}),
         ("   via the Lasserre hierarchy with localizing matrices:", {"size": 18})],
    ])
    add_text(s, Inches(0.6), Inches(3.8), Inches(12.1), Inches(0.6),
             "val^(k)(d)  ≤  val^(k+1)(d)  ≤  val(d)",
             size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_rich(s, Inches(0.6), Inches(4.65), Inches(12.1), Inches(0.6), [
        [("Step 3.  Certify", {"size": 18, "bold": True, "color": TEAL}),
         ("   a dual-feasible point of the order-k SDP whose objective exceeds the threshold.",
          {"size": 18})],
    ])

    add_rich(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.4), [
        [("Scaling lever", {"size": 16, "bold": True, "color": INK}),
         ("   —   correlative sparsity  (Waki–Kim–Kojima–Muramatsu 2006).",
          {"size": 16})],
        [("Chordal extension of the banded window graph yields  O(d)  overlapping cliques;",
          {"size": 15, "color": GREY})],
        [("the dense  binom(d+k, k)  moment matrix is replaced by  O(d)  small matrices of size  binom(b+1+k, k).",
          {"size": 15, "color": GREY})],
    ])

    add_footer(s, idx, total)


def slide_08_sdp_result(prs, idx, total):
    s = blank_slide(prs)
    slide_header(s, "Track 2 · Result",
                 "A certified lower bound with full rigor",
                 color=TEAL)

    s.shapes.add_picture(str(FIG / "fig_lasserre_ladder.png"),
                         Inches(0.4), Inches(1.8), width=Inches(7.5))

    add_box(s, Inches(8.1), Inches(1.85), Inches(4.7), Inches(2.3),
            fill=LIGHT, line=TEAL)
    add_text(s, Inches(8.3), Inches(1.95), Inches(4.3), Inches(0.4),
             "Certified", size=13, color=GREY)
    add_text(s, Inches(8.3), Inches(2.3), Inches(4.3), Inches(1.0),
             "C₁ₐ  ≥  1.3",
             size=36, bold=True, color=TEAL)
    add_text(s, Inches(8.3), Inches(3.3), Inches(4.3), Inches(0.5),
             "(d, k) = (16, 3)",
             size=14, color=INK)
    add_text(s, Inches(8.3), Inches(3.65), Inches(4.3), Inches(0.5),
             "MOSEK  ·  < 10 minutes",
             size=14, color=INK)

    add_rich(s, Inches(8.1), Inches(4.4), Inches(4.8), Inches(2.8), [
        [("What we record", {"size": 16, "bold": True, "color": TEAL})],
        [("• a single dual-feasible point  (y, Z, Λ)", {"size": 13})],
        [("• post-processed via Jansson–", {"size": 13})],
        [("   Chaykin–Keil interval-arithmetic shift", {"size": 13})],
        [("• auxiliary runs at  d ∈ {32, 64, 128}", {"size": 13})],
        [("   under correlative sparsity", {"size": 13})],
    ])

    add_footer(s, idx, total)


def slide_09_matlab_setup(prs, idx, total):
    s = blank_slide(prs)
    slide_header(s, "Track 3 · Auditing the C&S MATLAB",
                 "The code the authors shared with us — what it does",
                 color=CORAL)

    add_rich(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.8), [
        [("When we set out to reproduce the  C₁ₐ ≥ 1.2802  baseline we ran the MATLAB file",
          {"size": 17})],
        [("Cloninger and Steinerberger shared  (original_baseline_matlab.m).  The published",
          {"size": 17})],
        [("C&S paper is ", {"size": 17}),
         ("correct",  {"size": 17, "bold": True}),
         (" — it uses the fine height grid  B", {"size": 17}),
         ("n,m",  {"size": 12}),
         ("  with  m = 50.", {"size": 17})],
        [("", {"size": 6})],
        [("The MATLAB file, though, parameterizes the search differently:", {"size": 17})],
    ])

    add_box(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(1.8),
            fill=LIGHT, line=CORAL)
    add_rich(s, Inches(1.1), Inches(4.55), Inches(11.1), Inches(1.7), [
        [("In the paper", {"size": 16, "bold": True, "color": NAVY}),
         (":   step-function ", {"size": 16}),
         ("heights",  {"size": 16, "bold": True}),
         ("  quantized in steps of  ", {"size": 16}),
         ("1/m  =  0.02", {"size": 16, "bold": True})],
        [("", {"size": 4})],
        [("In the code", {"size": 16, "bold": True, "color": CORAL}),
         (":     bin ", {"size": 16}),
         ("masses", {"size": 16, "bold": True}),
         ("  quantized in steps of  ", {"size": 16}),
         ("0.02", {"size": 16, "bold": True}),
         ("   → heights quantized in steps of  ", {"size": 16}),
         ("2d × 0.02", {"size": 16, "bold": True, "color": CORAL})],
    ])

    add_text(s, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.8),
             "Different parameterization. Different height step. Same correction formula.",
             size=17, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    add_footer(s, idx, total)


def slide_10_matlab_gap(prs, idx, total):
    s = blank_slide(prs)
    slide_header(s, "The gap",
                 "Lemma 3 needs the height step; the code plugs in the mass step",
                 color=CORAL)

    add_rich(s, Inches(0.6), Inches(1.7), Inches(6.3), Inches(2.8), [
        [("C&S Lemma 3  gives", {"size": 16, "bold": True, "color": INK})],
        [("(g ∗ g)(x)  ≤  (a ∗ a)(x)  +  2ε  +  ε²",
          {"size": 17, "color": NAVY})],
        [("where  ε  is a bound on  ‖g − a‖∞  (a ", {"size": 15}),
         ("height",  {"size": 15, "italic": True}),
         ("  difference).",  {"size": 15})],
        [("", {"size": 6})],
        [("In the MATLAB code", {"size": 16, "bold": True, "color": CORAL})],
        [("ε  =  0.02  (mass step), but heights  h", {"size": 15}),
         ("i",  {"size": 11}),
         ("  =  2d · mass", {"size": 15}),
         ("i",  {"size": 11}),
         (",", {"size": 15})],
        [("so the true  ‖g − a‖∞  ≤  2d × 0.02,  not  0.02.", {"size": 15})],
        [("", {"size": 6})],
        [("The code's pruning budget is too small by a factor of  ≈ 2d.",
          {"size": 16, "bold": True, "color": CORAL})],
    ])

    s.shapes.add_picture(str(FIG / "fig_correction_gap.png"),
                         Inches(7.0), Inches(1.8), width=Inches(6.0))

    add_footer(s, idx, total)


def slide_11_smoking_gun(prs, idx, total):
    s = blank_slide(prs)
    slide_header(s, "The consequence",
                 "The code would \"prove away\" a known construction",
                 color=CORAL)

    add_rich(s, Inches(0.6), Inches(1.7), Inches(6.4), Inches(5.0), [
        [("The check", {"size": 18, "bold": True, "color": INK})],
        [("Matolcsi–Vinuesa  (2010)  exhibit a continuous  f*  with",
          {"size": 15})],
        [("‖f* ∗ f*‖∞  ≤  1.5029.", {"size": 15})],
        [("So  C₁ₐ  ≤  1.5029.", {"size": 15, "italic": True})],
        [("", {"size": 8})],
        [("Point the MATLAB cascade at the threshold  1.51:",
          {"size": 15})],
        [("• interior-window budget stays at ≈ 0.04 across cascade levels",
          {"size": 14})],
        [("• the step-function approximation of  f*  has self-convolution peaks",
          {"size": 14})],
        [("  that exceed 1.55 on narrow windows at moderate  d",
          {"size": 14})],
        [("• the code therefore prunes every descendant of  f*", {"size": 14})],
        [("• exhausted survivors → claim:  C₁ₐ  ≥  1.51", {"size": 14})],
        [("", {"size": 6})],
        [("But  1.51  >  1.5029  ≥  C₁ₐ.  Contradiction.",
          {"size": 16, "bold": True, "color": CORAL})],
    ])

    s.shapes.add_picture(str(FIG / "fig_effective_m.png"),
                         Inches(7.3), Inches(1.85), width=Inches(5.8))
    add_text(s, Inches(7.3), Inches(5.85), Inches(5.8), Inches(1.2),
             "Equivalently: the effective fine-grid  m  in C&S units is  1 / (2d ε)  — "
             "it drops below 1 by  d ≈ 25, at which point Lemma 3 says nothing at all.",
             size=12, color=GREY, italic=True)

    add_footer(s, idx, total)


def slide_12_what_this_means(prs, idx, total):
    s = blank_slide(prs)
    slide_header(s, "What we're saying — and what we're not",
                 "The published 1.2802 stands; the shared MATLAB file does not prove it",
                 color=CORAL)

    add_box(s, Inches(0.6), Inches(1.8), Inches(5.9), Inches(4.8),
            fill=LIGHT, line=NAVY)
    add_text(s, Inches(0.8), Inches(1.95), Inches(5.5), Inches(0.5),
             "The paper is fine", size=18, bold=True, color=NAVY)
    add_rich(s, Inches(0.8), Inches(2.45), Inches(5.5), Inches(4.0), [
        [("C&S 2017 work on the fine height grid  B", {"size": 14}),
         ("n,m", {"size": 10}),
         (",  where Lemma 3's", {"size": 14})],
        [("2/m + 1/m²  correction is legitimate.  The published bound", {"size": 14})],
        [("1.2802  is not in doubt.", {"size": 14})],
        [("", {"size": 8})],
        [("The MATLAB artifact is a different animal:", {"size": 14})],
        [("it enumerates a coarser mass grid, but ports the", {"size": 14})],
        [("paper's correction formula verbatim.  That combination", {"size": 14})],
        [("cannot certify the paper's bound, and, as shown, would", {"size": 14})],
        [("certify bounds that contradict the MV upper bound.", {"size": 14})],
    ])

    add_box(s, Inches(6.8), Inches(1.8), Inches(5.9), Inches(4.8),
            fill=LIGHT, line=CORAL)
    add_text(s, Inches(7.0), Inches(1.95), Inches(5.5), Inches(0.5),
             "Two live hypotheses", size=18, bold=True, color=CORAL)
    add_rich(s, Inches(7.0), Inches(2.45), Inches(5.5), Inches(4.0), [
        [("(a)  They sent us an early prototype.", {"size": 14, "bold": True})],
        [("The published bound was run on a separate fine-grid", {"size": 14})],
        [("enumerator; the coarse-mass file is a speed-oriented", {"size": 14})],
        [("variant never meant to be the proof vehicle.", {"size": 14})],
        [("", {"size": 8})],
        [("(b)  A derivation we're missing.", {"size": 14, "bold": True})],
        [("A mass-space bound that avoids the  2d  factor —", {"size": 14})],
        [("unlikely but worth asking about before we conclude.", {"size": 14})],
        [("", {"size": 8})],
        [("Either way, this is a narrow claim about one file.", {"size": 14, "italic": True, "color": GREY})],
    ])

    add_text(s, Inches(0.6), Inches(6.8), Inches(12.1), Inches(0.5),
             "We fixed the analogous bug in our own Python port on 2026-04-07 — "
             "switched to the fine grid, which is what makes the cascade's 1.4 sound.",
             size=13, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    add_footer(s, idx, total)


def slide_13_summary(prs, idx, total):
    s = blank_slide(prs)
    slide_header(s, "Takeaways",
                 "Where C₁ₐ stands after this project")

    s.shapes.add_picture(str(FIG / "fig_bounds_number_line.png"),
                         Inches(1.9), Inches(1.65), width=Inches(9.5))

    add_rich(s, Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.5), [
        [("Three deliverables", {"size": 20, "bold": True, "color": NAVY})],
        [("", {"size": 4})],
        [("1.  ", {"size": 16, "bold": True, "color": NAVY}),
         ("Cascade on GPU", {"size": 16, "bold": True}),
         ("  —  ", {"size": 16}),
         ("C₁ₐ ≥ 1.4", {"size": 16, "bold": True, "color": NAVY}),
         ("  via novel pruning bounds + fused CUDA kernels.",
          {"size": 16})],
        [("2.  ", {"size": 16, "bold": True, "color": TEAL}),
         ("Lasserre SDP certificate", {"size": 16, "bold": True}),
         ("  —  ", {"size": 16}),
         ("C₁ₐ ≥ 1.3 (certified)", {"size": 16, "bold": True, "color": TEAL}),
         ("  at  (d,k)=(16,3)  in under 10 minutes.",
          {"size": 16})],
        [("3.  ", {"size": 16, "bold": True, "color": CORAL}),
         ("MATLAB audit", {"size": 16, "bold": True}),
         ("  —  the shared C&S artifact confuses the mass step for the height step; ",
          {"size": 16}),
         ("the published paper is unaffected.",
          {"size": 16, "italic": True})],
        [("", {"size": 4})],
        [("Next:  formal write-ups for both proofs in ", {"size": 15, "color": GREY}),
         ("proof/", {"size": 15, "font": "Consolas", "color": GREY}),
         (",  push cascade to  d ≥ 256,  close the gap toward 1.5029.",
          {"size": 15, "color": GREY})],
    ])

    add_footer(s, idx, total)


def build():
    prs = new_prs()
    slide_01_title(prs)
    TOTAL = 13

    slide_02_problem(prs, 2, TOTAL)
    slide_03_prior_work(prs, 3, TOTAL)
    slide_04_three_tracks(prs, 4, TOTAL)
    slide_05_cascade_method(prs, 5, TOTAL)
    slide_06_cascade_result(prs, 6, TOTAL)
    slide_07_sdp_method(prs, 7, TOTAL)
    slide_08_sdp_result(prs, 8, TOTAL)
    slide_09_matlab_setup(prs, 9, TOTAL)
    slide_10_matlab_gap(prs, 10, TOTAL)
    slide_11_smoking_gun(prs, 11, TOTAL)
    slide_12_what_this_means(prs, 12, TOTAL)
    slide_13_summary(prs, 13, TOTAL)

    prs.save(OUT)
    print("wrote", OUT)


if __name__ == "__main__":
    build()
